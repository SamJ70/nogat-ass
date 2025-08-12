#!/usr/bin/env python3
"""
pptx_inconsistency_checker_final.py

Final working single-file CLI tool for analyzing .pptx decks for factual/logical inconsistencies.

Main design:
- Per-slide deterministic heuristics (fast, offline)
- Per-slide multimodal fact extraction using Gemini (text + images) -> structured facts JSON
- Deterministic comparisons of extracted facts to flag conflicts (numeric mismatch, percent sums, timeline mismatch, contradictory claims)
- Overlapping batching + character-based splitting to avoid token limits
- Robust fallbacks: if Gemini not available, heuristics + OCR still run
- Configurable CLI flags for cost control (--no-llm, --fact-only)

Requirements:
    pip install python-pptx pillow pytesseract google-generativeai python-dateutil fuzzywuzzy python-Levenshtein
    tesseract-ocr system binary installed (or configure pytesseract.pytesseract.tesseract_cmd)
    set GEMINI_API_KEY or GOOGLE_API_KEY or GENAI_API_KEY env var to your Gemini API key
"""

import os
import sys
import json
import tempfile
import re
import argparse
import traceback
from collections import defaultdict
from datetime import datetime
from dateutil import parser as dateparser
from math import isfinite, fabs
from typing import List, Dict, Any

from pptx import Presentation
from PIL import Image
import pytesseract

import google.generativeai as genai
from fuzzywuzzy import fuzz

# ------------- CONFIG -------------
GEMINI_MODEL = "gemini-2.5-flash"
MAX_SLIDES_PER_GEMINI_REQUEST = 10   # slides per batch
BATCH_OVERLAP = 3                    # overlap
OCR_CHUNK_LIMIT = 4000               # per-slide truncation for prompts
MAX_BATCH_CHARS = 80_000             # total chars per LLM call
NUMERIC_FUZZY_THRESHOLD = 60         # for grouping evidence in heuristics
METRIC_CANONICAL_FUZZY = 70          # for matching metrics across slides
MIN_RELEVANT_NUMBER = 1e-6           # avoid nonsense
PERCENT_TOLERANCE = 3.0              # percent-sum tolerance (absolute)
VERBOSE_DEFAULT = True
# -----------------------------------

# regex helpers
NUM_RE = re.compile(r'(?<!\w)(?:\$|€|£)?\s*([+-]?\d{1,3}(?:[,\d]{0,})?(?:\.\d+)?)(\s*%|\s*(?:million|bn|billion|k|M|B))?', re.IGNORECASE)
PERCENT_RE = re.compile(r'([+-]?\d+(?:\.\d+)?)\s*%')
DATE_TOKEN_RE = re.compile(r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Sept|Oct|Nov|Dec|\d{1,2})[^\n,./-]{0,20}\d{2,4}\b', re.IGNORECASE)

# global caches
_UPLOADED_IMAGE_CACHE: Dict[str, Any] = {}  # local path -> uploaded handle/object

# ----------------- Utilities -----------------
def safe_float(x):
    try:
        f = float(x)
        if not isfinite(f):
            return None
        return f
    except Exception:
        return None

def clean_text(s: str) -> str:
    return (s or "").strip()

def text_for_shape(shape):
    text = ""
    try:
        if shape.has_text_frame:
            text = "\n".join([p.text.strip() for p in shape.text_frame.paragraphs if p.text and p.text.strip()])
    except Exception:
        text = ""
    return text

# ----------------- PPTX parsing -----------------
def parse_pptx(path: str):
    prs = Presentation(path)
    slides = []
    tmpdir = tempfile.mkdtemp(prefix="pptx_imgs_")
    for i, slide in enumerate(prs.slides, start=1):
        stext = []
        tables = []
        saved_images = []
        for shape in slide.shapes:
            t = text_for_shape(shape)
            if t:
                stext.append(t)
            # tables
            try:
                if shape.has_table:
                    rows = []
                    for r in shape.table.rows:
                        rows.append([c.text.strip() for c in r.cells])
                    tables.append(rows)
            except Exception:
                pass
            # images
            try:
                if getattr(shape, "image", None) is not None:
                    img = shape.image
                    blob = img.blob
                    ext = img.ext or "png"
                    fname = os.path.join(tmpdir, f"slide{i}_img_{len(saved_images)}.{ext}")
                    with open(fname, "wb") as f:
                        f.write(blob)
                    saved_images.append(fname)
            except Exception:
                pass
        slides.append({
            "slide_index": i,
            "text": "\n".join(stext),
            "tables": tables,
            "images": saved_images
        })
    return slides

# ----------------- OCR -----------------
def ocr_image(path: str) -> str:
    try:
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        return text or ""
    except Exception:
        return ""

# ----------------- Heuristic Detectors (fast) -----------------
def _is_metric_word_nearby(segment: str) -> bool:
    """Return True if segment contains metric keywords making the number meaningful."""
    kws = ["revenue", "sales", "profit", "loss", "hours", "hrs", "minute", "min", "month", "annual", "year", "per", "avg", "users", "customers", "employees", "people", "headcount", "kpi", "$", "€", "£", "%"]
    seg = segment.lower()
    for kw in kws:
        if kw in seg:
            return True
    return False

def detect_numeric_conflicts(slides: List[Dict[str, Any]], fuzzy_threshold=NUMERIC_FUZZY_THRESHOLD):
    mentions = []
    window = 40
    for s in slides:
        txt = s.get("text", "") or ""
        segments = [seg.strip() for seg in re.split(r'[\r\n]+', txt) if seg.strip()]
        for seg in segments:
            for m in NUM_RE.finditer(seg):
                start, end = m.span()
                left = seg[max(0, start - window):start].strip()
                right = seg[end:end + window].strip()
                snippet = (left + " " + m.group(0) + " " + right).strip()
                # heuristics: discard trivial numbers unless a metric word nearby
                if not _is_metric_word_nearby(left + " " + right):
                    # allow percentages even if no metric word nearby
                    if not PERCENT_RE.search(m.group(0)):
                        # also allow large numbers with currency symbols
                        if not re.search(r'[\$\€\£]', m.group(0)):
                            # skip trivial numbers
                            continue
                left_words = left.split()[-6:] if left else []
                right_words = right.split()[:6] if right else []
                key = re.sub(r'\d+[\d,\.]*', '', ' '.join(left_words + right_words)).strip()
                val = safe_float(m.group(1).replace(",", ""))
                mentions.append({
                    "slide": s.get("slide_index"),
                    "snippet": snippet,
                    "value": val,
                    "raw": m.group(0),
                    "key": key
                })

    # grouping
    groups = []
    used = [False] * len(mentions)
    for i, a in enumerate(mentions):
        if used[i]:
            continue
        grp = [a]
        used[i] = True
        for j in range(i + 1, len(mentions)):
            if used[j]:
                continue
            score = 0
            if a.get("key") and mentions[j].get("key"):
                score = fuzz.partial_ratio(a["key"][:80], mentions[j]["key"][:80])
            else:
                score = fuzz.partial_ratio(a["snippet"][:80], mentions[j]["snippet"][:80])
            if score > fuzzy_threshold:
                grp.append(mentions[j])
                used[j] = True
        if len(grp) > 1:
            vals = [g["value"] for g in grp if g["value"] is not None]
            if len(vals) > 1 and len({round(v, 6) for v in vals}) > 1:
                slides_ref = sorted(list({x["slide"] for x in grp if x["slide"] is not None}))
                groups.append({"group": grp, "slides": slides_ref})
    issues = []
    for g in groups:
        issues.append({
            "type": "numeric_conflict",
            "slides": g["slides"],
            "evidence": g["group"]
        })
    return issues

def detect_percentage_sum_problems(slides: List[Dict[str, Any]], tolerance=PERCENT_TOLERANCE):
    issues = []
    for s in slides:
        txt = s.get("text", "") or ""
        if re.search(r'\b(breakdown|split|composition|distribution|share|mix|percentages?)\b', txt, re.IGNORECASE):
            percents = [safe_float(m.group(1)) for m in PERCENT_RE.finditer(txt)]
            percents = [p for p in percents if p is not None]
            if percents and abs(sum(percents) - 100.0) > tolerance:
                issues.append({"type": "percent_sum_mismatch", "slides": [s["slide_index"]], "sum": sum(percents), "percents": percents, "text_sample": txt[:300]})
    return issues

def detect_timeline_mismatches(slides: List[Dict[str, Any]]):
    slide_dates = {}
    for s in slides:
        dts = extract_dates_from_text(s.get("text", "") or "")
        slide_dates[s["slide_index"]] = dts
    issues = []
    n = len(slides)
    for i in range(n):
        for j in range(i + 1, n):
            di = slide_dates.get(i + 1, [])
            dj = slide_dates.get(j + 1, [])
            if di and dj:
                try:
                    di0 = dateparser.parse(di[0]["parsed"])
                    dj0 = dateparser.parse(dj[0]["parsed"])
                    if di0 > dj0 and (j + 1) > (i + 1):
                        issues.append({"type": "timeline_mismatch", "slides": [i + 1, j + 1], "dates": [di[0], dj[0]], "note": "earlier slide contains a later date than a later slide"})
                except Exception:
                    pass
    return issues

def extract_dates_from_text(text: str):
    dates = []
    for m in DATE_TOKEN_RE.finditer(text):
        token = m.group(0)
        try:
            dt = dateparser.parse(token, default=datetime(2000, 1, 1))
            if dt:
                dates.append({"raw": token, "parsed": dt.isoformat()})
        except Exception:
            continue
    iso_dates = re.findall(r'\b\d{4}-\d{2}-\d{2}\b', text)
    for d in iso_dates:
        try:
            dt = dateparser.parse(d)
            dates.append({"raw": d, "parsed": dt.isoformat()})
        except Exception:
            pass
    return dates

# ----------------- Gemini client helpers -----------------
def init_gemini_client():
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY") or os.getenv("GENAI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY (or GOOGLE_API_KEY / GENAI_API_KEY) environment variable not set.")
    genai.configure(api_key=api_key)
    try:
        model = genai.GenerativeModel(GEMINI_MODEL)
    except Exception:
        model = None
    return model

def _call_model_generate_content(client, contents):
    """
    Unified call to model.generate_content with fallback for SDK shapes.
    `client` may be a GenerativeModel instance or None.
    `contents` is a list of parts (dict/text) or a single string.
    """
    model_obj = client or genai.GenerativeModel(GEMINI_MODEL)
    try:
        # prefer named arg
        return model_obj.generate_content(content=contents)
    except TypeError:
        try:
            return model_obj.generate_content(contents)
        except Exception as e:
            raise

def _extract_genai_text(response):
    if isinstance(response, str):
        return response
    if hasattr(response, "text") and response.text:
        return response.text
    try:
        out = getattr(response, "output", None)
        if out:
            if isinstance(out, (list, tuple)) and len(out) > 0:
                c0 = out[0]
                try:
                    return c0["content"][0]["text"]
                except Exception:
                    try:
                        return c0.content[0].text
                    except Exception:
                        pass
    except Exception:
        pass
    if hasattr(response, "candidates"):
        try:
            cand = response.candidates[0]
            if hasattr(cand, "content"):
                return cand.content[0].text if cand.content else str(cand)
        except Exception:
            pass
    return str(response)

def try_extract_json_from_text(text: str):
    """Try to find a JSON array/object in model text; return parsed or None."""
    if not text:
        return None
    m = re.search(r'(\[\s*\{.*\}\s*\])', text, re.S) or re.search(r'(\[.*\])', text, re.S)
    if not m:
        # sometimes model prints JSON-like lines; try to clean full text
        try:
            return json.loads(text)
        except Exception:
            return None
    candidate = m.group(1)
    try:
        return json.loads(candidate)
    except Exception:
        try:
            t = candidate.replace("'", '"')
            t = re.sub(r',\s*}', '}', t)
            t = re.sub(r',\s*\]', ']', t)
            return json.loads(t)
        except Exception:
            return None

# ----------------- Fact extraction & canonicalization -----------------
def canonicalize_metric_name(name: str) -> str:
    """Lightweight canonicalization: lowercase, remove punctuation, remove stopwords, collapse spaces."""
    if not name:
        return ""
    s = name.lower()
    s = re.sub(r'[^a-z0-9\s]', ' ', s)
    tokens = [t for t in s.split() if t not in {"the", "a", "per", "of", "in", "for", "and", "to", "by", "on", "monthly", "annual", "yearly"}]
    canon = " ".join(tokens).strip()
    return re.sub(r'\s+', ' ', canon)

def merge_similar_metrics(facts: List[Dict[str, Any]]):
    """
    Cluster metrics by canonical name + fuzzy match, produce canonical_key for each fact.
    Returns list of facts with 'canon' added.
    """
    canon_map = []
    for f in facts:
        metric = canonicalize_metric_name(f.get("metric") or f.get("name") or "")
        matched = None
        for c in canon_map:
            if fuzz.token_set_ratio(metric, c) >= METRIC_CANONICAL_FUZZY:
                matched = c
                break
        if matched is None:
            canon_map.append(metric)
            f["canon"] = metric
        else:
            f["canon"] = matched
    return facts

# ----------------- LLM per-slide fact extraction (multimodal) -----------------
def build_fact_extraction_prompt(slide: Dict[str, Any]) -> List[Any]:
    """
    Build a small `contents` list for generate_content that includes:
      - brief instruction
      - slide text (truncated)
      - attached image uploads or OCR snippets
    Returns a list to pass to model.generate_content
    """
    instr = (
        "Extract structured facts from the provided slide. "
        "Return a JSON array where each entry is an object: "
        '{"metric": "<short name>", "value": <number|null>, "unit": "<unit or % or currency>", '
        '"qualifier": "<e.g. per consultant, monthly, annual, forecast/baseline>", "excerpt": "<short exact excerpt from slide>", "slide": <slide_index>}. '
        "If no numeric facts present, return an empty array []. "
        "Do not add extra commentary outside the JSON array."
    )
    contents = [{"text": instr}]
    text = slide.get("text", "") or ""
    if len(text) > OCR_CHUNK_LIMIT:
        text = text[:OCR_CHUNK_LIMIT] + "\n... [truncated]"
    contents.append({"text": f"Slide {slide['slide_index']} TEXT:\n\"\"\"\n{text}\n\"\"\""})
    # attach images (uploaded) or OCR fallback
    for p in slide.get("images", []):
        uploaded = _UPLOADED_IMAGE_CACHE.get(p)
        if uploaded is None and hasattr(genai, "upload_file"):
            try:
                uploaded = genai.upload_file(p)
            except Exception:
                uploaded = None
            _UPLOADED_IMAGE_CACHE[p] = uploaded
        if uploaded is not None:
            # many SDKs accept the uploaded object directly in contents
            contents.append(uploaded)
        else:
            ocr_snip = ocr_image(p)[:800]
            contents.append({"text": f"[IMAGE FILE: {os.path.basename(p)}] OCR (first 800 chars): {ocr_snip}"})
    return contents

def extract_facts_for_slides(client, slides: List[Dict[str, Any]], verbose=True, only_slide_indices: List[int] = None):
    """
    For each slide (or only selected indices), run a per-slide multimodal extraction to produce structured facts.
    Returns list of fact dicts.
    """
    facts = []
    client = client  # may be None; _call_model_generate_content will instantiate if needed
    for s in slides:
        idx = s["slide_index"]
        if only_slide_indices and idx not in only_slide_indices:
            continue
        if verbose:
            print(f"  extracting facts from slide {idx}...")
        try:
            contents = build_fact_extraction_prompt(s)
            # guard by max chars
            total_chars = sum(len(c.get("text", "")) if isinstance(c, dict) and "text" in c else 0 for c in contents)
            if total_chars > MAX_BATCH_CHARS:
                # truncate slide text more aggressively
                for part in contents:
                    if isinstance(part, dict) and "text" in part and len(part["text"]) > 2000:
                        part["text"] = part["text"][:2000] + "\n... [truncated]"
            resp = _call_model_generate_content(client, contents)
            raw = _extract_genai_text(resp)
            parsed = try_extract_json_from_text(raw)
            if parsed is None:
                # fallback: log raw output as lrm_raw_output fact for later human triage
                facts.append({"metric": "__llm_raw__", "value": None, "unit": None, "qualifier": None, "excerpt": raw[:500], "slide": idx, "llm_raw": raw})
            else:
                # ensure each fact has keys and slide
                for item in parsed:
                    if isinstance(item, dict):
                        item.setdefault("metric", item.get("name") or "")
                        item.setdefault("value", item.get("value", None))
                        item.setdefault("unit", item.get("unit", None))
                        item.setdefault("qualifier", item.get("qualifier", None))
                        item.setdefault("excerpt", item.get("excerpt", "")[:400])
                        item.setdefault("slide", idx)
                        # convert value to number if possible
                        if isinstance(item["value"], str):
                            v = re.sub(r'[^\d\.\-]', '', item["value"])
                            item["value"] = safe_float(v)
                        facts.append(item)
        except Exception as e:
            facts.append({"metric": "__llm_error__", "value": None, "unit": None, "qualifier": None, "excerpt": str(e), "slide": idx})
    # canonicalize/merge metric names
    facts = merge_similar_metrics(facts)
    return facts

# ----------------- Deterministic comparisons using extracted facts -----------------
def compare_facts_and_flag(facts: List[Dict[str, Any]]):
    """
    Using structured facts, flag:
     - numeric_conflict: same canon metric present on multiple slides with different values (thresholded)
     - percent_sum_mismatch: within a slide, percentages not summing to ~100 (we also have heuristic for these)
     - timeline_mismatch: dates/qualifiers contradictions inferred from qualifiers and dates
    """
    issues = []
    # index facts by canon metric
    by_canon = defaultdict(list)
    for f in facts:
        canon = f.get("canon") or canonicalize_metric_name(f.get("metric") or "")
        by_canon[canon].append(f)

    # numeric conflicts
    for canon, flist in by_canon.items():
        vals = [f.get("value") for f in flist if isinstance(f.get("value"), (int, float))]
        if len(vals) > 1:
            # relative threshold: consider conflict if not approximately equal
            unique_vals = {round(v, 6) for v in vals}
            if len(unique_vals) > 1:
                slides_ref = sorted(list({f["slide"] for f in flist if f.get("slide")}))
                issues.append({"type": "numeric_conflict", "metric": canon, "slides": slides_ref, "evidence": flist})

    # percent-sum per slide (from facts)
    by_slide = defaultdict(list)
    for f in facts:
        by_slide[f.get("slide")].append(f)
    for slide_idx, flist in by_slide.items():
        percents = [f.get("value") for f in flist if f.get("unit") and isinstance(f.get("value"), (int, float)) and '%' in (f.get("unit") or "")]
        if percents and abs(sum(percents) - 100.0) > PERCENT_TOLERANCE:
            issues.append({"type": "percent_sum_mismatch", "slides": [slide_idx], "sum": sum(percents), "percents": percents, "evidence": flist})

    # timeline mismatch (simple): compare facts with qualifier containing 'by'/'forecast' vs 'baseline' and dates present in excerpt
    # This is heuristic — more advanced tagging could improve this.
    dated_facts = []
    for f in facts:
        excerpt = f.get("excerpt") or ""
        dts = extract_dates_from_text(excerpt)
        if dts:
            f["_dates"] = dts
            dated_facts.append(f)
    # compare pairs
    for i in range(len(dated_facts)):
        for j in range(i + 1, len(dated_facts)):
            fi = dated_facts[i]
            fj = dated_facts[j]
            try:
                di = dateparser.parse(fi["_dates"][0]["parsed"])
                dj = dateparser.parse(fj["_dates"][0]["parsed"])
                # if earlier slide mentions later date than later slide -> mismatch
                if fi.get("slide") < fj.get("slide") and di > dj:
                    issues.append({"type": "timeline_mismatch", "slides": [fi["slide"], fj["slide"]], "dates": [fi["_dates"][0], fj["_dates"][0]], "evidence": [fi, fj]})
            except Exception:
                pass

    return issues

# ----------------- batching helpers -----------------
def overlapping_batches(items: List[Any], size=MAX_SLIDES_PER_GEMINI_REQUEST, overlap=BATCH_OVERLAP):
    if size <= 0:
        yield items
        return
    n = len(items)
    if n == 0:
        return
    i = 0
    step = size - overlap if size > overlap else size
    yielded = set()
    while i < n:
        batch = items[i:i + size]
        key = tuple([s["slide_index"] for s in batch])
        if key not in yielded:
            yielded.add(key)
            yield batch
        i += step
        if i + size >= n and i < n:
            final_batch = items[max(0, n - size):n]
            key2 = tuple([s["slide_index"] for s in final_batch])
            if key2 not in yielded:
                yield final_batch
            break

def split_slide_blurbs_by_chars(slide_blurbs: List[Dict[str, Any]], max_chars=MAX_BATCH_CHARS):
    windows = []
    current = []
    current_len = 0
    for s in slide_blurbs:
        est = len(s.get("text", "")) + 200
        if current and (current_len + est > max_chars):
            windows.append(current)
            current = []
            current_len = 0
        current.append(s)
        current_len += est
    if current:
        windows.append(current)
    return windows

# ----------------- LLM chunk-based contradiction detection (fallback) -----------------
def ask_gemini_contradictions_chunk(client, slides_chunk: List[Dict[str, Any]]):
    """
    If you still want a chunk-based 'ask LLM to find contradictions across a chunk',
    this function performs a robust multi-modal call and returns parsed JSON or raw output wrapped.
    This is kept as a fallback; primary comparisons use per-slide fact extraction + deterministic comparisons.
    """
    instructions = (
        "You are an assistant that compares slide texts and associated images to find factual/logical inconsistencies across them.\n"
        "Return a JSON array where each item has: type, slides (list), explanation, excerpts (map slide->excerpt).\n"
        "Types: numeric_conflict, percent_mismatch, timeline_mismatch, contradictory_claim, other.\n"
        "If none, return []."
    )
    # build contents list
    contents = [{"text": instructions}]
    for s in slides_chunk:
        t = s.get("text", "") or ""
        if len(t) > OCR_CHUNK_LIMIT:
            t = t[:OCR_CHUNK_LIMIT] + "\n... [truncated]"
        contents.append({"text": f"Slide {s['slide_index']} TEXT:\n\"\"\"\n{t}\n\"\"\""})
        for p in s.get("images", []):
            uploaded = _UPLOADED_IMAGE_CACHE.get(p)
            if uploaded is None and hasattr(genai, "upload_file"):
                try:
                    uploaded = genai.upload_file(p)
                except Exception:
                    uploaded = None
                _UPLOADED_IMAGE_CACHE[p] = uploaded
            if uploaded is not None:
                contents.append(uploaded)
            else:
                contents.append({"text": f"[IMAGE FILE: {os.path.basename(p)}] OCR: {ocr_image(p)[:400]}"})
    try:
        resp = _call_model_generate_content(client, contents)
        raw = _extract_genai_text(resp)
        parsed = try_extract_json_from_text(raw)
        if parsed is None:
            return [{"type": "llm_raw_output", "slides": [s["slide_index"] for s in slides_chunk], "explanation": raw}]
        return parsed
    except Exception as e:
        return [{"type": "llm_error", "error": str(e), "slides": [s["slide_index"] for s in slides_chunk], "raw": traceback.format_exc()}]

# ----------------- Orchestration -----------------
def analyze_deck(pptx_path: str, extra_images_dir: str = None, no_llm: bool = False, fact_only: bool = False, verbose: bool = VERBOSE_DEFAULT):
    slides = parse_pptx(pptx_path)

    # append OCR text from images so heuristics have some visual text
    for s in slides:
        for imgpath in s.get("images", []):
            try:
                ocr_text = ocr_image(imgpath)
                if ocr_text and ocr_text.strip():
                    s["text"] = (s.get("text", "") or "") + "\n\n[OCR IMAGE TEXT]\n" + ocr_text
            except Exception:
                pass

    # additional external images as pseudo-slides
    if extra_images_dir and os.path.isdir(extra_images_dir):
        for i, fname in enumerate(sorted(os.listdir(extra_images_dir)), start=len(slides) + 1):
            path = os.path.join(extra_images_dir, fname)
            if os.path.isfile(path):
                ocr_text = ocr_image(path) or ""
                slides.append({"slide_index": i, "text": f"[external image {fname}]\n{ocr_text}", "tables": [], "images": [path]})

    report = {"analyzed_at": datetime.utcnow().isoformat() + "Z", "source": pptx_path, "issues": [], "facts": []}

    # run fast heuristics
    if verbose:
        print("Running deterministic heuristics...")
    try:
        report["issues"].extend(detect_numeric_conflicts(slides))
        report["issues"].extend(detect_percentage_sum_problems(slides))
        report["issues"].extend(detect_timeline_mismatches(slides))
    except Exception as e:
        report["issues"].append({"type": "heuristics_error", "error": str(e)})

    if no_llm:
        if verbose:
            print("Skipping LLM (Gemini) steps (--no-llm).")
        return report

    # init client
    try:
        client = init_gemini_client()
    except Exception as e:
        client = None
        if verbose:
            print("Warning: Gemini init failed:", e)
        # do not fail — still run heuristics & OCR

    # Per-slide fact extraction: (recommended approach)
    if verbose:
        print("Extracting structured facts (per-slide multimodal fact extraction)...")
    # To reduce cost: only extract facts for slides that heuristics flagged OR all slides depending on strategy.
    flagged_slides = set()
    for it in report["issues"]:
        for s in it.get("slides", []) or []:
            flagged_slides.add(s)
    # Strategy: extract facts for flagged slides plus a sample of other slides (here we extract for all slides to be thorough).
    # You can change to: only_slide_indices = sorted(flagged_slides) to reduce cost.
    only_slide_indices = None  # set to None to extract for all slides; or set to sorted(flagged_slides) for cheaper run
    facts = extract_facts_for_slides(client, slides, verbose=verbose, only_slide_indices=only_slide_indices)
    report["facts"] = facts

    # deterministic comparisons over facts
    if verbose:
        print("Comparing extracted facts to find deterministic conflicts...")
    try:
        fact_issues = compare_facts_and_flag(facts)
        report["issues"].extend(fact_issues)
    except Exception as e:
        report["issues"].append({"type": "fact_compare_error", "error": str(e)})

    # fallback: chunk-based LLM contradiction detection (optional, expensive)
    if not fact_only:
        if verbose:
            print("Running chunk-based LLM contradiction detection (fallback, may use API quota)...")
        batches = list(overlapping_batches(slides, size=MAX_SLIDES_PER_GEMINI_REQUEST, overlap=BATCH_OVERLAP))
        seen_sets = set()
        unique_batches = []
        for b in batches:
            idxs = tuple([s["slide_index"] for s in b])
            if idxs in seen_sets:
                continue
            seen_sets.add(idxs)
            unique_batches.append(b)
        for batch in unique_batches:
            # further split by chars to avoid token limits
            slide_blurbs = [{"slide_index": s["slide_index"], "text": (s.get("text","") or "")[:OCR_CHUNK_LIMIT], "images": s.get("images", [])} for s in batch]
            windows = split_slide_blurbs_by_chars(slide_blurbs, max_chars=MAX_BATCH_CHARS)
            for win in windows:
                lres = ask_gemini_contradictions_chunk(client, win)
                if isinstance(lres, list):
                    for item in lres:
                        if isinstance(item, dict):
                            if "slide" in item and "slides" not in item:
                                item["slides"] = item.get("slide")
                            report["issues"].append(item)
                        else:
                            report["issues"].append({"type": "llm_unexpected_item", "content": item, "slides": [s["slide_index"] for s in win]})
                else:
                    report["issues"].append({"type": "llm_unexpected", "content": lres, "slides": [s["slide_index"] for s in win]})

    # coarse de-dup / canonical merging of issues
    merged = []
    seen = set()
    for it in report["issues"]:
        try:
            # Normalize key for dedupe (skip ephemeral fields)
            keyobj = {"type": it.get("type"), "slides": it.get("slides"), "metric": it.get("metric")}
            key = json.dumps(keyobj, sort_keys=True, default=str)
        except Exception:
            key = str(it)
        if key in seen:
            # try to merge evidence if present
            for m in merged:
                try:
                    mk = json.dumps({"type": m.get("type"), "slides": m.get("slides"), "metric": m.get("metric")}, sort_keys=True, default=str)
                except Exception:
                    mk = ""
                if mk == key:
                    # merge evidence arrays
                    if "evidence" in it:
                        m.setdefault("evidence", [])
                        # append items that are not already in evidence (coarse)
                        for ev in it.get("evidence", []):
                            if ev not in m["evidence"]:
                                m["evidence"].append(ev)
                    break
            continue
        seen.add(key)
        merged.append(it)
    report["issues"] = merged

    return report

# ----------------- CLI -----------------
def main():
    p = argparse.ArgumentParser(description="Detect inconsistencies in a PPTX deck (hybrid heuristics + Gemini multimodal facts).")
    p.add_argument("pptx", help="path to .pptx")
    p.add_argument("--images", help="optional directory of slide images to OCR", default=None)
    p.add_argument("--out", help="output JSON file", default="inconsistency_report_final.json")
    p.add_argument("--no-llm", help="skip any LLM/Gemini calls (run heuristics + OCR only)", action="store_true")
    p.add_argument("--fact-only", help="skip chunk-based LLM contradictions; rely only on fact-extraction comparisons", action="store_true")
    p.add_argument("--quiet", help="less console output", action="store_true")
    args = p.parse_args()

    if not os.path.isfile(args.pptx):
        print("pptx not found:", args.pptx)
        sys.exit(2)

    verbose = not args.quiet
    if verbose:
        print("Parsing PPTX and running heuristics + (optional) Gemini multimodal checks ...")

    report = analyze_deck(args.pptx, args.images, no_llm=args.no_llm, fact_only=args.fact_only, verbose=verbose)

    with open(args.out, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2, ensure_ascii=False)
    print("Done. Report written to", args.out)

    # summary
    print("\nSummary of issues found:")
    if not report.get("issues"):
        print("No issues detected.")
        return
    for idx, issue in enumerate(report["issues"], start=1):
        itype = issue.get("type")
        slides = issue.get("slides") or issue.get("slide") or issue.get("slide_index")
        explanation = (issue.get("explanation") or issue.get("note") or issue.get("error") or issue.get("metric") or "")[:200]
        print(f"{idx}. {itype} — slides {slides} — {explanation}")

if __name__ == "__main__":
    main()
