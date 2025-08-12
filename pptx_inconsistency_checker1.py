#!/usr/bin/env python3
"""
pptx_inconsistency_checker_multimodal.py

Implements the feedback:
 - pass images (not just OCR) to the LLM for spatial context
 - move from brittle rules -> LLM-driven fact extraction (multimodal) + deterministic checks
 - perform global comparisons over compact facts so distant contradictions are found

Usage:
    python pptx_inconsistency_checker_multimodal.py /path/to/deck.pptx [--images /path/to/images_dir] [--out report.json] [--no-llm]

Requirements:
    pip install python-pptx pillow pytesseract google-generativeai python-dateutil fuzzywuzzy python-Levenshtein

Environment:
    Set GEMINI_API_KEY or GOOGLE_API_KEY or GENAI_API_KEY to a Gemini API key.

Notes:
 - The script is defensive for different genai SDK shapes.
 - Per-slide fact extraction is multimodal: each slide's text + uploaded image(s) (if upload supported) are sent to Gemini.
 - Global comparisons are done on extracted facts (compact), but the model is also asked to cross-check facts globally.
"""

import os
import sys
import json
import tempfile
import re
import argparse
import traceback
from datetime import datetime
from typing import List, Dict, Any
from math import isfinite

from pptx import Presentation
from PIL import Image
import pytesseract

import google.generativeai as genai  # ensure installed
from fuzzywuzzy import fuzz
from dateutil import parser as dateparser

# ---------------- CONFIG ----------------
GEMINI_MODEL = "gemini-2.5-flash"
OCR_CHUNK_LIMIT = 4000
MAX_FACTS_JSON_CHARS = 150_000  # keep global facts payload reasonable
UPLOAD_CACHE: Dict[str, Any] = {}  # path -> uploaded handle/object
# ----------------------------------------

# Regex helpers
NUM_RE = re.compile(r'(?<!\w)(?:\$|€|£)?\s*([+-]?\d{1,3}(?:[,\d]{0,})?(?:\.\d+)?)(\s*%|\s*(?:million|bn|billion|k|M|B))?', re.IGNORECASE)
PERCENT_RE = re.compile(r'([+-]?\d+(?:\.\d+)?)\s*%')

# ---------------- Utilities ----------------
def safe_float(x):
    try:
        f = float(x)
        if not isfinite(f):
            return None
        return f
    except Exception:
        return None

def text_for_shape(shape):
    text = ""
    try:
        if shape.has_text_frame:
            text = "\n".join([p.text.strip() for p in shape.text_frame.paragraphs if p.text and p.text.strip()])
    except Exception:
        text = ""
    return text

# ---------------- PPTX parsing ----------------
def parse_pptx(path: str) -> List[Dict[str, Any]]:
    prs = Presentation(path)
    slides = []
    tmpdir = tempfile.mkdtemp(prefix="pptx_imgs_")
    for i, slide in enumerate(prs.slides, start=1):
        stext = []
        tables = []
        imgs = []
        for shape in slide.shapes:
            t = text_for_shape(shape)
            if t:
                stext.append(t)
            try:
                if shape.has_table:
                    rows = []
                    for r in shape.table.rows:
                        rows.append([c.text.strip() for c in r.cells])
                    tables.append(rows)
            except Exception:
                pass
            try:
                if getattr(shape, "image", None) is not None:
                    img = shape.image
                    blob = img.blob
                    ext = img.ext or "png"
                    fname = os.path.join(tmpdir, f"slide{i}_img_{len(imgs)}.{ext}")
                    with open(fname, "wb") as f:
                        f.write(blob)
                    imgs.append(fname)
            except Exception:
                pass
        slides.append({
            "slide_index": i,
            "text": "\n".join(stext),
            "tables": tables,
            "images": imgs
        })
    return slides

# ---------------- OCR fallback ----------------
def ocr_image(path: str) -> str:
    try:
        img = Image.open(path)
        return pytesseract.image_to_string(img) or ""
    except Exception:
        return ""

# ---------------- Gemini client ----------------
def init_genai():
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("GOOGLE_API_KEY") or os.getenv("GENAI_API_KEY")
    if not api_key:
        raise RuntimeError("Set GEMINI_API_KEY / GOOGLE_API_KEY / GENAI_API_KEY in env.")
    genai.configure(api_key=api_key)
    # return model instance or None (we'll create later if needed)
    try:
        model = genai.GenerativeModel(GEMINI_MODEL)
    except Exception:
        model = None
    return model

def _call_generate_content(model_obj, contents):
    """
    Unified call to generate_content; handles slight SDK differences.
    contents: either a string or a list of parts (dicts or uploaded handles).
    """
    model = model_obj or genai.GenerativeModel(GEMINI_MODEL)
    try:
        # prefer named arg
        return model.generate_content(content=contents)
    except TypeError:
        return model.generate_content(contents)

def _extract_text_from_response(resp) -> str:
    if isinstance(resp, str):
        return resp
    if hasattr(resp, "text") and resp.text:
        return resp.text
    try:
        out = getattr(resp, "output", None)
        if out:
            if isinstance(out, (list, tuple)) and len(out) > 0:
                c0 = out[0]
                try:
                    return c0["content"][0]["text"]
                except Exception:
                    try:
                        return c0.content[0].text
                    except Exception:
                        pass
    except Exception:
        pass
    if hasattr(resp, "candidates"):
        try:
            cand = resp.candidates[0]
            if hasattr(cand, "content"):
                return cand.content[0].text if cand.content else str(cand)
        except Exception:
            pass
    return str(resp)

def try_parse_json_from_text(text: str):
    if not text:
        return None
    m = re.search(r'(\[\s*\{.*\}\s*\])', text, re.S) or re.search(r'(\[.*\])', text, re.S)
    if m:
        cand = m.group(1)
        try:
            return json.loads(cand)
        except Exception:
            try:
                t = cand.replace("'", '"')
                t = re.sub(r',\s*}', '}', t)
                t = re.sub(r',\s*\]', ']', t)
                return json.loads(t)
            except Exception:
                return None
    # try entire text
    try:
        return json.loads(text)
    except Exception:
        return None

# ---------------- Upload helper ----------------
def upload_image_if_supported(path: str):
    # reuse cache
    if path in UPLOAD_CACHE:
        return UPLOAD_CACHE[path]
    handle = None
    if hasattr(genai, "upload_file"):
        try:
            handle = genai.upload_file(path)
        except Exception:
            handle = None
    UPLOAD_CACHE[path] = handle
    return handle

# ---------------- Fact extraction (multimodal per-slide) ----------------
def build_fact_extraction_contents(slide: Dict[str, Any]) -> List[Any]:
    """
    Build the 'contents' list for model.generate_content where we include:
     - instruction text
     - slide text part
     - attached uploaded image object(s) when available
     - if upload not available, include OCR snippet as text
    """
    instruction = (
        "Extract a JSON array of **numeric facts** from the slide.\n"
        "Each fact object should contain:\n"
        '  - metric: short string name (e.g., "Revenue", "Hours saved per consultant")\n'
        '  - value: numeric value or null\n'
        '  - unit: unit string (e.g., "%", "USD", "hours", "people") or null\n'
        '  - qualifier: short qualifier (e.g., "monthly", "annual", "forecast", "baseline") or null\n'
        '  - excerpt: exact short excerpt (<=200 chars) from the slide or image that supports the fact\n'
        'Return ONLY the JSON array (no commentary). If no numeric facts, return []\n'
    )
    parts = [{"text": instruction}]
    text = slide.get("text", "") or ""
    if len(text) > OCR_CHUNK_LIMIT:
        text = text[:OCR_CHUNK_LIMIT] + "\n... [truncated]"
    parts.append({"text": f"Slide {slide['slide_index']} TEXT:\n\"\"\"\n{text}\n\"\"\""})
    # images: try to attach uploaded handle if SDK supports it
    for imgpath in slide.get("images", []):
        handle = upload_image_if_supported(imgpath)
        if handle is not None:
            parts.append(handle)
        else:
            # include OCR fallback snippet
            snip = ocr_image(imgpath)[:800]
            parts.append({"text": f"[IMAGE FILE: {os.path.basename(imgpath)}] OCR (first 800 chars): {snip}"})
    return parts

def extract_facts_multimodal(model_obj, slides: List[Dict[str, Any]], verbose: bool = True) -> List[Dict[str, Any]]:
    """
    For each slide, call Gemini with multimodal content (attached image handles when available)
    to extract structured facts. Returns a flat list of facts with slide index.
    """
    facts = []
    for s in slides:
        idx = s["slide_index"]
        if verbose:
            print(f"Extracting facts from slide {idx} (multimodal)...")
        try:
            contents = build_fact_extraction_contents(s)
            resp = _call_generate_content(model_obj, contents)
            raw = _extract_text_from_response(resp)
            parsed = try_parse_json_from_text(raw)
            if parsed is None:
                # fallback: attempt to extract numbers with heuristics if LLM returns non-json
                facts.append({"metric": "__llm_raw__", "value": None, "unit": None, "qualifier": None, "excerpt": raw[:300], "slide": idx})
            else:
                for item in parsed:
                    if not isinstance(item, dict):
                        continue
                    item.setdefault("metric", item.get("metric") or item.get("name") or "")
                    # try coerce numeric
                    v = item.get("value")
                    if isinstance(v, str):
                        vnum = re.sub(r'[^\d\.\-]', '', v)
                        item["value"] = safe_float(vnum)
                    elif isinstance(v, (int, float)):
                        item["value"] = float(v)
                    else:
                        item["value"] = None
                    item.setdefault("unit", item.get("unit"))
                    item.setdefault("qualifier", item.get("qualifier"))
                    item.setdefault("excerpt", (item.get("excerpt") or "")[:400])
                    item["slide"] = idx
                    facts.append(item)
        except Exception as e:
            facts.append({"metric": "__llm_error__", "value": None, "unit": None, "qualifier": None, "excerpt": str(e), "slide": idx})
    return facts

# ---------------- Deterministic checks on facts ----------------
def canonicalize_metric(name: str) -> str:
    if not name:
        return ""
    s = name.lower()
    s = re.sub(r'[^a-z0-9\s]', ' ', s)
    s = re.sub(r'\b(per|the|a|an|of|in|for|and|to|by)\b', ' ', s)
    s = re.sub(r'\s+', ' ', s).strip()
    return s

def compare_facts_deterministic(facts: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    issues = []
    # group by canonical metric
    groups = {}
    for f in facts:
        canon = canonicalize_metric(f.get("metric") or "")
        groups.setdefault(canon, []).append(f)
    # numeric conflict: same canon metric across slides, differing numeric values
    for canon, fl in groups.items():
        vals = [f["value"] for f in fl if isinstance(f.get("value"), (int, float))]
        if len(vals) > 1:
            unique = {round(v,6) for v in vals}
            if len(unique) > 1:
                slides = sorted({f["slide"] for f in fl})
                issues.append({"type": "numeric_conflict", "metric": canon or None, "slides": slides, "evidence": fl})
    # percent-sum: within slide
    by_slide = {}
    for f in facts:
        by_slide.setdefault(f["slide"], []).append(f)
    for slide, fl in by_slide.items():
        percents = [f["value"] for f in fl if isinstance(f.get("value"), (int,float)) and (f.get("unit") and "%" in f.get("unit"))]
        if percents and abs(sum(percents) - 100.0) > 3.0:
            issues.append({"type": "percent_sum_mismatch", "slides": [slide], "sum": sum(percents), "percents": percents, "evidence": fl})
    # timeline mismatches - basic: look for dates in excerpts and compare ordering vs slide order
    dated = []
    for f in facts:
        ex = (f.get("excerpt") or "")
        dates = re.findall(r'\b\d{4}\b', ex)  # naive year extract
        if dates:
            try:
                # pick first year
                year = int(dates[0])
                dated.append((f, year))
            except Exception:
                pass
    for i in range(len(dated)):
        for j in range(i+1, len(dated)):
            fi, yi = dated[i]
            fj, yj = dated[j]
            if fi["slide"] < fj["slide"] and yi > yj:
                issues.append({"type": "timeline_mismatch", "slides": [fi["slide"], fj["slide"]], "dates": [yi, yj], "evidence": [fi, fj]})
    return issues

# ---------------- Global LLM cross-check on compact facts ----------------
def ask_global_contradictions_on_facts(model_obj, facts: List[Dict[str, Any]], verbose=True):
    """
    Give the LLM the compact JSON facts (already multimodal-extracted) and ask for contradictions across the deck.
    This keeps token size small because facts are compact; LLM will produce high-level contradictions using the richer visual-aware facts.
    """
    # build compact facts JSON string (truncated/resized if needed)
    facts_payload = json.dumps(facts, default=str)
    if len(facts_payload) > MAX_FACTS_JSON_CHARS:
        # trim raw excerpt sizes to shrink payload and reserialize
        small = []
        for f in facts:
            f2 = dict(f)
            if "excerpt" in f2 and isinstance(f2["excerpt"], str):
                f2["excerpt"] = f2["excerpt"][:200]
            small.append(f2)
        facts_payload = json.dumps(small, default=str)
    instruction = (
        "You are given a JSON array of extracted facts from slides. Each fact has: metric, value (number|null), unit, qualifier, excerpt, slide.\n"
        "Compare across the facts and return a JSON array of inconsistencies found. Each inconsistency should be an object with fields:\n"
        " - type (one of numeric_conflict, percent_sum_mismatch, timeline_mismatch, contradictory_claim, other)\n"
        " - slides: list of slide indices involved\n"
        " - explanation: short human-readable explanation\n"
        " - evidence: map slide->excerpt or a list of excerpts\n        \n"
        "Return ONLY the JSON array.\n"
    )
    contents = [{"text": instruction}, {"text": facts_payload}]
    try:
        resp = _call_generate_content(model_obj, contents)
        raw = _extract_text_from_response(resp)
        parsed = try_parse_json_from_text(raw)
        if parsed is None:
            # return raw as llm_raw_output so human triage can inspect
            return [{"type": "llm_raw_output", "explanation": raw[:1000]}]
        return parsed
    except Exception as e:
        return [{"type": "llm_error", "error": str(e), "raw": traceback.format_exc()}]

# ---------------- Orchestration ----------------
def analyze_deck(pptx_path: str, images_dir: str = None, out_json: str = "report.json", no_llm: bool = False, verbose: bool = True):
    slides = parse_pptx(pptx_path)

    # attach OCR text for heuristics if images exist
    for s in slides:
        for img in s.get("images", []):
            try:
                ocrt = ocr_image(img)
                if ocrt and ocrt.strip():
                    s["text"] = (s.get("text","") or "") + "\n\n[OCR IMAGE TEXT]\n" + ocrt
            except Exception:
                pass

    # add external images as pseudo-slides if provided
    if images_dir and os.path.isdir(images_dir):
        for i, fname in enumerate(sorted(os.listdir(images_dir)), start=len(slides)+1):
            path = os.path.join(images_dir, fname)
            if os.path.isfile(path):
                slides.append({"slide_index": i, "text": f"[external image {fname}]\n", "tables": [], "images": [path]})

    report = {
        "analyzed_at": datetime.utcnow().isoformat() + "Z",
        "source": pptx_path,
        "issues": [],
        "facts": []
    }

    # quick heuristics (keeps behavior if no LLM)
    # Numeric heuristic: detect obvious repeated/conflicting numbers (conservative)
    if verbose:
        print("Running conservative numeric heuristics (fallback if LLM unavailable)...")
    # conservative mention extracting: only numbers with currency, percent or metric words nearby
    mentions = []
    for s in slides:
        txt = s.get("text","") or ""
        for m in NUM_RE.finditer(txt):
            raw = m.group(0)
            # require currency or percent or metric words nearby (simple heuristic)
            context = txt[max(0, m.start()-40): m.end()+40].lower()
            if ('$' in raw) or ('€' in raw) or ('£' in raw) or ('%' in raw) or any(k in context for k in ["revenue","sales","profit","hours","monthly","annual","users","customers","per"]):
                val = safe_float(m.group(1).replace(",",""))
                mentions.append({"slide": s["slide_index"], "raw": raw, "value": val, "context": context})
    # group crude conflicts by same context fuzz
    # (kept intentionally conservative; main detection uses LLM below)
    used = [False]*len(mentions)
    for i,a in enumerate(mentions):
        if used[i]: continue
        grp=[a]; used[i]=True
        for j in range(i+1,len(mentions)):
            if used[j]: continue
            score = fuzz.partial_ratio(a['context'][:80], mentions[j]['context'][:80])
            if score>70:
                grp.append(mentions[j]); used[j]=True
        vals=[g['value'] for g in grp if g['value'] is not None]
        if len(vals)>1 and len({round(v,6) for v in vals})>1:
            report['issues'].append({"type":"numeric_conflict","slides": sorted(list({g['slide'] for g in grp})), "evidence": grp})

    # if LLM disabled, finish and return heuristics-only report
    if no_llm:
        if verbose:
            print("LLM disabled (--no-llm). Returning heuristics-only report.")
        return report

    # init model
    try:
        model = init_genai()
    except Exception as e:
        model = None
        if verbose:
            print("Failed to init Gemini client:", str(e))
            print("Proceeding with heuristics-only output.")
        return report

    # PER-SLIDE MULTIMODAL FACT EXTRACTION (images attached directly when upload supported)
    if verbose:
        print("Running per-slide multimodal fact extraction (images passed to LLM)...")
    facts = extract_facts_multimodal(model, slides, verbose=verbose)
    report["facts"] = facts

    # Deterministic comparisons on facts
    if verbose:
        print("Running deterministic comparisons on extracted facts...")
    try:
        det_issues = compare_facts_deterministic(facts)
        report["issues"].extend(det_issues)
    except Exception as e:
        report["issues"].append({"type":"compare_error","error":str(e)})

    # GLOBAL LLM-BASED CROSS-CHECK on compact facts
    if verbose:
        print("Asking LLM to cross-check facts globally (compact facts payload)...")
    try:
        global_llm_issues = ask_global_contradictions_on_facts(model, facts, verbose=verbose)
        # normalize and append
        if isinstance(global_llm_issues, list):
            for it in global_llm_issues:
                report["issues"].append(it)
        else:
            report["issues"].append({"type":"llm_unexpected","content":global_llm_issues})
    except Exception as e:
        report["issues"].append({"type":"global_llm_error","error":str(e)})

    # final coarse dedupe of issues
    seen=set(); dedup=[]
    for it in report["issues"]:
        try:
            key=json.dumps({"type":it.get("type"), "slides": it.get("slides"), "metric": it.get("metric")}, sort_keys=True, default=str)
        except Exception:
            key=str(it)
        if key in seen:
            continue
        seen.add(key)
        dedup.append(it)
    report["issues"]=dedup

    return report

# ---------------- CLI ----------------
def main():
    p = argparse.ArgumentParser(description="Multi-modal PPTX inconsistency checker (images passed to Gemini).")
    p.add_argument("pptx", help="path to .pptx")
    p.add_argument("--images", help="optional extra images dir", default=None)
    p.add_argument("--out", help="output JSON file", default="inconsistency_report_multimodal.json")
    p.add_argument("--no-llm", help="skip any LLM calls (heuristics+OCR only)", action="store_true")
    p.add_argument("--quiet", help="less console output", action="store_true")
    args = p.parse_args()

    if not os.path.isfile(args.pptx):
        print("pptx not found:", args.pptx); sys.exit(2)
    verbose = not args.quiet

    report = analyze_deck(args.pptx, images_dir=args.images, out_json=args.out, no_llm=args.no_llm, verbose=verbose)

    with open(args.out, "w", encoding="utf-8") as f:
        json.dump(report, f, indent=2, ensure_ascii=False)
    print("Wrote report to", args.out)
    print("\nSummary:")
    if not report.get("issues"):
        print("No issues found.")
        return
    for i, it in enumerate(report["issues"], start=1):
        typ = it.get("type")
        slides = it.get("slides") or it.get("slide") or "?"
        expl = it.get("explanation") or it.get("metric") or it.get("error") or ""
        print(f"{i}. {typ} — slides {slides} — {str(expl)[:200]}")

if __name__ == "__main__":
    # small alias to UPLOAD_CACHE for inner functions
    global UPLOAD_CACHE
    UPLOAD_CACHE = UPLOAD_CACHE  # type: ignore
    main()
